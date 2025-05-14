#!/usr/bin/python
# -*- coding: utf-8 -*-


"""
[env]
# Conda Environment
conda create --name ia_using_faiss python=3.9.13
conda info --envs
source activate ia_using_faiss
source activate clip_env
conda deactivate


# BURN AFTER READING
source activate ia_using_faiss

# if needed to remove
conda env remove -n [NAME_OF_THE_CONDA_ENVIRONMENT]
conda env remove -n ia_using_faiss


# install packages with conda
conda install -c conda-forge sentence-transformers
conda install -c pytorch faiss-cpu

# install packages with pip
python -m pip install sentence-transformers
python -m pip install pytorch faiss-cpu
python -m pip install numpy



# update conda 
conda update -n base -c defaults conda

# to export requirements
pip freeze > requirements.txt

# [path]
cd /Users/brunoflaven/Documents/01_work/blog_articles/ia_generate_slides/


# launch the file
python 004_pptx_python_script.py




"""

from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Introduction
slide_1 = prs.slides.add_slide(prs.slide_layouts[1])
title_1 = slide_1.shapes.title
content_1 = slide_1.placeholders[1]
title_1.text = "Introduction"
content_1.text = (
    "Définition du changement climatique : Le changement climatique désigne les variations à long terme des températures et des conditions météorologiques. "
    "Ces changements peuvent être naturels ou causés par les activités humaines."
)

# Slide 2: Causes du changement climatique
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
content_2 = slide_2.placeholders[1]
title_2.text = "Causes du changement climatique"
content_2.text = (
    "Gaz à effet de serre : Les principaux gaz à effet de serre sont le dioxyde de carbone (CO2), le méthane (CH4) et le protoxyde d'azote (N2O). "
    "Ces gaz piègent la chaleur dans l'atmosphère, ce qui entraîne une augmentation des températures.\n\n"
    "Activités humaines : La combustion de combustibles fossiles, la déforestation et certaines pratiques agricoles contribuent à l'augmentation des gaz à effet de serre."
)

# Slide 3: Effets du changement climatique
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
content_3 = slide_3.placeholders[1]
title_3.text = "Effets du changement climatique"
content_3.text = (
    "Températures mondiales : Les températures moyennes ont augmenté de 2°C en Europe contre 1,1°C en moyenne globale.\n\n"
    "Événements météorologiques extrêmes : L'intensité et le nombre de phénomènes météorologiques extrêmes, tels que les sécheresses, les vagues de chaleur, les pluies intenses, les inondations, les vents violents et les feux de forêt, ont augmenté.\n\n"
    "Impacts sur la santé : Le changement climatique a des effets négatifs sur la santé humaine, les infrastructures, l'énergie, les ressources en eau et l'économie."
)

# Slide 4: Solutions pour lutter contre le changement climatique
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
content_4 = slide_4.placeholders[1]
title_4.text = "Solutions pour lutter contre le changement climatique"
content_4.text = (
    "Atténuation : Réduire les émissions de gaz à effet de serre en adoptant des pratiques durables, telles que l'utilisation de sources d'énergie renouvelables et la réduction de la consommation d'énergie.\n\n"
    "Adaptation : Mettre en place des mesures pour s'adapter aux effets du changement climatique, telles que la gestion des ressources en eau et la mobilisation des connaissances des populations locales."
)

# Slide 5: Conclusion
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
content_5 = slide_5.placeholders[1]
title_5.text = "Conclusion"
content_5.text = (
    "Importance de l'action : Il est crucial de prendre des mesures pour lutter contre le changement climatique et s'adapter à ses effets. "
    "Une approche coordonnée, systémique et inclusive est nécessaire pour réduire les risques de manière efficace."
)

# Slide 6: Références
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
content_6 = slide_6.placeholders[1]
title_6.text = "Références"
content_6.text = (
    "Ministères Aménagement du territoire Transition écologique : Changement climatique : causes, effets et enjeux.\n\n"
    "GIEC : 10 points clés du rapport AR6 WGII.\n\n"
    "Gouvernement du Québec : Impacts des changements climatiques."
)

# Save the presentation to a file
prs.save("002_impact_du_changement_climatique.pptx")

print("La présentation PowerPoint a été créée avec succès.")


