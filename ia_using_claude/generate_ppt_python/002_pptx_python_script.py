#!/usr/bin/python
# -*- coding: utf-8 -*-


"""
[env]
# Conda Environment
conda create --name ia_using_faiss python=3.9.13
conda info --envs
source activate ia_using_faiss
source activate clip_env
conda deactivate


# BURN AFTER READING
source activate ia_using_faiss

# if needed to remove
conda env remove -n [NAME_OF_THE_CONDA_ENVIRONMENT]
conda env remove -n ia_using_faiss


# install packages with conda
conda install -c conda-forge sentence-transformers
conda install -c pytorch faiss-cpu

# install packages with pip
python -m pip install sentence-transformers
python -m pip install pytorch faiss-cpu
python -m pip install numpy



# update conda 
conda update -n base -c defaults conda

# to export requirements
pip freeze > requirements.txt

# [path]
cd /Users/brunoflaven/Documents/01_work/blog_articles/ia_generate_slides/


# launch the file
python 002_pptx_python_script.py




"""


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Charger la présentation existante
prs = Presentation('impact_du_changement_climatique.pptx')

# Définir une fonction pour appliquer un modèle à chaque diapositive
def apply_template(slide):
    # Définir la couleur de fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # Fond blanc

    # Définir le style du titre
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(24)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)  # Couleur bleu foncé
                paragraph.alignment = PP_ALIGN.CENTER

    # Définir le style du contenu
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)
                    run.font.color.rgb = RGBColor(0, 0, 0)  # Couleur noire
                paragraph.alignment = PP_ALIGN.LEFT

# Appliquer le modèle à chaque diapositive de la présentation
for slide in prs.slides:
    apply_template(slide)

# Enregistrer la présentation améliorée
prs.save('impact_du_changement_climatique_v1.pptx')

print("La présentation améliorée a été enregistrée sous le nom 'impact_du_changement_climatique_v1.pptx'.")